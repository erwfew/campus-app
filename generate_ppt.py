from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
BG_DARK = RGBColor(0x0a, 0x0a, 0x1a)
BG_CARD = RGBColor(0x16, 0x16, 0x2e)
ACCENT = RGBColor(0x43, 0x61, 0xee)
PURPLE = RGBColor(0xa7, 0x8b, 0xfa)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x99, 0x99, 0xbb)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x88)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, icon, title, desc, tag=''):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    shape.shadow.inherit = False

    add_text(slide, left + 0.3, top + 0.2, width - 0.6, 0.5, icon, font_size=28)
    add_text(slide, left + 0.3, top + 0.7, width - 0.6, 0.4, title, font_size=18, bold=True)
    add_text(slide, left + 0.3, top + 1.15, width - 0.6, height - 1.5, desc, font_size=12, color=GRAY)
    if tag:
        add_text(slide, left + 0.3, top + height - 0.5, 1.2, 0.3, tag, font_size=10, color=ACCENT)

def add_highlight(slide, left, top, width, icon, title, desc):
    add_text(slide, left, top, 0.6, 0.5, icon, font_size=24)
    add_text(slide, left + 0.7, top, width - 0.7, 0.35, title, font_size=16, bold=True)
    add_text(slide, left + 0.7, top + 0.4, width - 0.7, 0.5, desc, font_size=11, color=GRAY)

# ==================== 第1页：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, BG_DARK)
add_text(slide, 3, 1.5, 7.3, 1, '🎓', font_size=60, align=PP_ALIGN.CENTER)
add_text(slide, 2, 2.8, 9.3, 1, '校园综合服务', font_size=48, bold=True, align=PP_ALIGN.CENTER,
         color=WHITE)
add_text(slide, 2, 3.9, 9.3, 0.6, '一个 App，搞定你的校园日常', font_size=22, color=GRAY, align=PP_ALIGN.CENTER)

tags = ['📱 uni-app 跨平台', '🗺️ 智能导航', '📚 课程管理', '🏃 运动追踪', '🛒 文创商城']
tag_start = 3.2
for i, tag in enumerate(tags):
    add_text(slide, tag_start + i * 1.5, 5.0, 1.4, 0.4, tag, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ==================== 第2页：产品概述 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '01 · 产品概述', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '为什么需要校园综合服务？', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 8, 0.4, '一站式解决大学生校园生活中的高频需求', font_size=16, color=GRAY)

cards = [
    ('🎯', '痛点分析', '课表、导航、跑步、活动需要切换多个 App，信息分散，体验割裂'),
    ('💡', '解决方案', '将课程表、校园导航、运动追踪、文创商城、排行榜整合到一个平台'),
    ('👥', '目标用户', '在校大学生，尤其是新生——快速熟悉校园、高效管理学习生活'),
    ('✨', '核心价值', '减少 App 切换，提升校园生活效率，增加校园归属感和参与感'),
]
positions = [(0.8, 2.4), (4.6, 2.4), (0.8, 4.6), (4.6, 4.6)]
for (icon, title, desc), (x, y) in zip(cards, positions):
    add_card(slide, x, y, 3.5, 1.9, icon, title, desc)

# ==================== 第3页：功能总览 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '02 · 功能架构', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '六大核心模块', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 8, 0.4, '覆盖学习、出行、运动、消费、社交全场景', font_size=16, color=GRAY)

modules = [
    ('🏠', '首页仪表盘', '问候语、吉祥物、快捷入口、今日课程、校园公告', '入口'),
    ('🗺️', '校园导航', 'GPS 定位、路线规划、语音导航、3D 地图', '核心功能'),
    ('📚', '课程中心', '三种课表视图、作业管理、出勤统计', '学习'),
    ('🏃', '校园跑步', 'GPS 轨迹、实时配速、卡路里、历史记录', '运动'),
    ('🛒', '文创商城', '校园文创产品、活动报名、分类筛选', '消费'),
    ('👤', '个人中心', '签到系统、排行榜、成就徽章、学信网认证', '社交'),
]
positions = [(0.8, 2.4), (5.0, 2.4), (9.2, 2.4), (0.8, 4.6), (5.0, 4.6), (9.2, 4.6)]
for (icon, title, desc, tag), (x, y) in zip(modules, positions):
    add_card(slide, x, y, 3.8, 1.9, icon, title, desc, tag)

# ==================== 第4页：首页 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '03 · 首页仪表盘', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '你的校园生活一览', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 10, 0.4, '智能问候 · 吉祥物互动 · 课程提醒 · 公告推送', font_size=16, color=GRAY)

items = [
    ('🐱', '虚拟伙伴"小喵"', '互动式吉祥物，点击触发对话，通过签到和使用积累经验值不断升级'),
    ('📅', '今日课程实时状态', '根据当前时间自动标记已完成、上课中、即将开始，无需手动刷新'),
    ('📢', '校园公告推送', '热门、最新、通知三类标签，点击查看详情，不错过校园重要信息'),
    ('⚡', '快捷入口', '一键直达导航、课表、跑步、文创四大核心功能，减少操作步骤'),
]
for i, (icon, title, desc) in enumerate(items):
    add_highlight(slide, 0.8, 2.5 + i * 1.15, 11.5, icon, title, desc)

# ==================== 第5页：校园导航 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '04 · 校园导航', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '智能校园地图导航', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 10, 0.4, 'GPS 定位 · 路线规划 · 语音播报 · 3D 视角', font_size=16, color=GRAY)

nav_cards = [
    ('📍', '精准定位', 'WGS84 坐标系实时 GPS 定位，自动转换为 GCJ02 火星坐标，14 个校园地标一键快定位'),
    ('🛤️', '智能路线规划', '支持起点/终点/途经点设置，步行和骑行双模式，高德 API 优先 + 本地模拟兜底'),
    ('🗣️', '语音导航', 'Android TTS 原生引擎 + H5 SpeechSynthesis 降级，上下文感知播报，解放双眼'),
    ('🎮', '3D 实时导航', '地图 40° 倾斜 3D 视角，8 方向箭头指示，每 3 秒刷新位置，50 米内自动到达'),
]
positions = [(0.8, 2.4), (4.6, 2.4), (0.8, 4.6), (4.6, 4.6)]
for (icon, title, desc), (x, y) in zip(nav_cards, positions):
    add_card(slide, x, y, 3.5, 1.9, icon, title, desc)

# ==================== 第6页：课程中心 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '05 · 课程中心', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '高效课程与作业管理', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 10, 0.4, '三种视图 · 灵活配置 · 作业追踪 · 出勤统计', font_size=16, color=GRAY)

course_cards = [
    ('📊', '三种课表视图', '表格视图看全周概览，周视图按天浏览，日视图聚焦今天，今日课程高亮显示'),
    ('⚙️', '灵活配置', '自定义学期日期、每天节数、每节课时长、上课时间、课间休息，适配任何学校安排'),
    ('📝', '作业通知', '教师发布作业自动同步，截止日期颜色标记紧急程度，一键标记完成，状态本地持久化'),
    ('📈', '出勤统计', '出勤率、已上课时、缺勤次数三维度统计，数据可视化一目了然'),
]
positions = [(0.8, 2.4), (4.6, 2.4), (0.8, 4.6), (4.6, 4.6)]
for (icon, title, desc), (x, y) in zip(course_cards, positions):
    add_card(slide, x, y, 3.5, 1.9, icon, title, desc)

# ==================== 第7页：校园跑步 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '06 · 校园跑步', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, 'GPS 实时跑步追踪', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 10, 0.4, '轨迹绘制 · 实时配速 · 卡路里计算 · 历史记录', font_size=16, color=GRAY)

run_items = [
    ('🗺️', '实时轨迹绘制', '每 3 秒采集 GPS 位置，Haversine 公式精确计算距离，自动过滤 50 米以上漂移数据'),
    ('⏱️', '实时数据面板', '计时器、距离、配速、卡路里四维数据实时更新，暂停/继续/停止三态控制'),
    ('📋', '智能历史记录', '最多保存 50 条记录，自动分类为晨跑/日间跑/夜跑，每次跑步生成总结报告'),
    ('🔋', '优雅降级', 'GPS 信号弱时自动切换为时间估算模式（0.003 km/s），确保功能始终可用'),
]
for i, (icon, title, desc) in enumerate(run_items):
    add_highlight(slide, 0.8, 2.5 + i * 1.15, 11.5, icon, title, desc)

# ==================== 第8页：文创商城 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '07 · 文创商城', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '校园文创与活动平台', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 10, 0.4, '特色商品 · 分类筛选 · 活动报名 · 社区互动', font_size=16, color=GRAY)

add_card(slide, 0.8, 2.4, 5.5, 3.0, '🛍️', '文创商品',
         '8 款校园特色文创产品：徽章笔记本、校园背包、钥匙扣、明信片、卫衣、帆布包、保温杯、贴纸包\n\n支持文具/饰品/服装分类筛选，底部弹窗查看详情，一键加入购物车')
add_card(slide, 6.8, 2.4, 5.5, 3.0, '🎉', '热门活动',
         '校园文化节、草地音乐节、读书分享会、摄影大赛等丰富活动\n\n查看活动时间地点详情，一键报名，状态实时更新为"已报名"')

# ==================== 第9页：个人中心 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '08 · 个人中心', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '数据汇聚与成就体系', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 10, 0.4, '每日签到 · 排行榜 · 成就徽章 · 身份认证', font_size=16, color=GRAY)

profile_cards = [
    ('✅', '每日签到', '每日签到获得 10 经验值，助力小喵升级，培养用户粘性'),
    ('🏆', '多维排行榜', '到课率、课堂表现、运动排行三个维度，支持班级/年级/学院三级范围筛选'),
    ('🎖️', '成就徽章', '学习之星、运动达人、全勤奖等 6+ 徽章，解锁条件可视化，记录校园荣耀'),
    ('🔐', '学信网认证', '嵌入 web-view 对接学信网，4 步动画认证流程，保障校园安全与真实性'),
]
positions = [(0.8, 2.4), (4.6, 2.4), (0.8, 4.6), (4.6, 4.6)]
for (icon, title, desc), (x, y) in zip(profile_cards, positions):
    add_card(slide, x, y, 3.5, 1.9, icon, title, desc)

# ==================== 第10页：技术架构 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 0.8, 0.5, 5, 0.4, '09 · 技术架构', font_size=13, color=ACCENT, bold=True)
add_text(slide, 0.8, 0.9, 8, 0.7, '技术选型与核心能力', font_size=36, bold=True)
add_text(slide, 0.8, 1.6, 10, 0.4, '跨平台框架 · 地图引擎 · 数据持久化 · iOS/Android 双端', font_size=16, color=GRAY)

# 三列技术卡片
for i, (icon, title, items) in enumerate([
    ('📱', '前端框架', ['uni-app 跨平台开发', 'Vue 2 + Options API', '毛玻璃卡片 UI 设计', 'iOS safe-area 适配', '原生 tabBar + 自定义导航栏']),
    ('🗺️', '地图与定位', ['腾讯地图 SDK', '高德地图路线 API', 'WGS84/GCJ02 坐标转换', 'Catmull-Rom 路线模拟', 'Android TTS 语音引擎']),
    ('💾', '数据与后端', ['Express.js RESTful API', 'localStorage 本地持久化', '学信网 web-view 对接', '学分/作业服务端同步', '10+ 数据存储 Key']),
]):
    x = 0.8 + i * 4.1
    shape = slide.shapes.add_shape(1, Inches(x), Inches(2.4), Inches(3.7), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()

    add_text(slide, x + 0.3, 2.6, 3.1, 0.5, icon, font_size=24)
    add_text(slide, x + 0.3, 3.1, 3.1, 0.4, title, font_size=18, bold=True, color=PURPLE)
    for j, item in enumerate(items):
        add_text(slide, x + 0.3, 3.7 + j * 0.55, 3.1, 0.45, '› ' + item, font_size=12, color=GRAY)

# ==================== 第11页：结尾 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, 3, 1.5, 7.3, 1, '🎓', font_size=56, align=PP_ALIGN.CENTER)
add_text(slide, 2, 2.8, 9.3, 1, '校园综合服务', font_size=48, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2, 3.9, 9.3, 0.6, '你的校园生活，一个 App 搞定', font_size=22, color=GRAY, align=PP_ALIGN.CENTER)

end_tags = ['📱 支持 iOS & Android', '🔧 uni-app 跨平台', '🚀 v1.0.0']
for i, tag in enumerate(end_tags):
    add_text(slide, 3.5 + i * 2.3, 5.2, 2.2, 0.5, tag, font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 保存
output_path = r'C:\Users\吕嘉成\AppData\Roaming\CherryStudio\Data\Agents\d8xgjfrd2\campus-app\校园综合服务-产品演示.pptx'
prs.save(output_path)
print(f'PPT 已生成: {output_path}')
